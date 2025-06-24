import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from pathlib import Path

def read_document_text(file_path: str) -> str:
    """
    Reads text from a document file (PDF, DOCX, PPTX).

    Args:
        file_path: The path to the document file.

    Returns:
        The extracted text content of the document.
    
    Raises:
        ValueError: If the file format is not supported.
    """
    path = Path(file_path)
    suffix = path.suffix.lower()

    text_content = ""

    if suffix == '.pdf':
        with fitz.open(path) as doc:
            for page in doc:
                text_content += page.get_text()
    elif suffix == '.docx':
        doc = Document(path)
        for para in doc.paragraphs:
            text_content += para.text + '\n'
    elif suffix == '.pptx':
        pres = Presentation(path)
        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content += shape.text + '\n'
    else:
        raise ValueError(f"Unsupported file format: {suffix}")

    return text_content.strip() 